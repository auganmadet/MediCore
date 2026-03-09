#!/usr/bin/env python3
"""Générateur de présentation PowerPoint MediCore aux couleurs Médiprix."""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Couleurs Médiprix
RED = RGBColor(0xC8, 0x10, 0x2E)
DARK_RED = RGBColor(0x8B, 0x0A, 0x1E)
LIGHT_RED = RGBColor(0xF2, 0xD5, 0xDA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x1A, 0x56, 0x8E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x6A, 0x00)

LOGO_PATH = "docs/logo-mediprix.png"
LOGOS_DIR = "docs/logos"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_counter = 0
TOTAL_SLIDES = 21


def next_slide_num():
    """Incrémente et retourne le numéro de slide courant."""
    global slide_counter
    slide_counter += 1
    return slide_counter


def logo_path(name):
    """Retourne le chemin du logo si le fichier existe, sinon None."""
    path = os.path.join(LOGOS_DIR, f"logo-{name}.png")
    return path if os.path.exists(path) else None


def safe_add_picture(slide, name, left, top, height):
    """Ajoute un logo s'il existe, sinon ne fait rien."""
    path = logo_path(name)
    if path:
        slide.shapes.add_picture(path, left, top, height=height)


def add_red_bar(slide, top=0, height=Inches(1.2)):
    """Barre rouge en haut du slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), top, prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide):
    """Barre fine rouge en bas du slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Inches(7.1), prs.slide_width, Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_RED
    shape.line.fill.background()
    return shape


def add_logo(slide, left=Inches(11.5), top=Inches(0.15), height=Inches(0.9)):
    """Logo Médiprix en haut à droite."""
    slide.shapes.add_picture(LOGO_PATH, left, top, height=height)


def add_page_number(slide, num):
    """Numéro de page en bas à droite."""
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.12), Inches(1.1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {TOTAL_SLIDES}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT


def _add_styled_bullet(tf, text, is_first, is_sub=False):
    """Ajoute un bullet avec le texte avant ':' en rouge gras."""
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()

    if is_sub:
        p.level = 1
        run = p.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY
        p.space_before = Pt(4)
    elif ':' in text and not text.startswith('  '):
        parts = text.split(':', 1)
        run1 = p.add_run()
        run1.text = parts[0] + " :"
        run1.font.size = Pt(18)
        run1.font.bold = True
        run1.font.color.rgb = RED

        run2 = p.add_run()
        run2.text = parts[1]
        run2.font.size = Pt(18)
        run2.font.color.rgb = BLACK
        p.space_before = Pt(12)
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(18)
        run.font.color.rgb = BLACK
        p.space_before = Pt(12)

    p.space_after = Pt(4)
    return p


def _add_styled_bullet_col(tf, text, is_first, font_size=16):
    """Bullet stylé pour colonnes (taille configurable)."""
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()

    if ':' in text:
        parts = text.split(':', 1)
        run1 = p.add_run()
        run1.text = parts[0] + " :"
        run1.font.size = Pt(font_size)
        run1.font.bold = True
        run1.font.color.rgb = RED

        run2 = p.add_run()
        run2.text = parts[1]
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = BLACK
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = BLACK

    p.space_before = Pt(8)
    p.space_after = Pt(2)
    return p


def _add_sommaire_line(tf, text, is_first):
    """Ligne du sommaire avec numéro et titre en rouge gras."""
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()

    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RED

    p.space_before = Pt(14)
    p.space_after = Pt(6)
    return p


def add_title_slide(title, subtitle):
    """Slide de titre (couverture)."""
    num = next_slide_num()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_red_bar(slide, top=Inches(2.2), height=Inches(3.0))
    pic = slide.shapes.add_picture(LOGO_PATH, Inches(0), Inches(0.5), height=Inches(1.4))
    pic.left = int((prs.slide_width - pic.width) / 2)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.3), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT_RED
    p2.alignment = PP_ALIGN.CENTER

    add_bottom_bar(slide)
    add_page_number(slide, num)
    return slide


def add_sommaire_slide(items):
    """Slide sommaire avec titres en rouge gras."""
    num = next_slide_num()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_red_bar(slide, top=0, height=Inches(1.2))
    add_bottom_bar(slide)
    add_page_number(slide, num)
    add_logo(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Sommaire"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.0), Inches(5.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(items):
        _add_sommaire_line(tf2, item, i == 0)
    return slide


def add_section_slide(section_num, section_title):
    """Slide de transition entre sections."""
    num = next_slide_num()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RED

    slide.shapes.add_picture(LOGO_PATH, Inches(11.3), Inches(0.3), height=Inches(0.8))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(2.3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(36)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    add_page_number(slide, num)
    return slide


def add_content_slide(title, bullets):
    """Slide de contenu avec termes clés en rouge gras."""
    num = next_slide_num()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_red_bar(slide, top=0, height=Inches(1.2))
    add_bottom_bar(slide)
    add_page_number(slide, num)
    add_logo(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        is_sub = bullet.startswith("  - ")
        text = bullet[4:] if is_sub else bullet
        _add_styled_bullet(tf2, text, i == 0, is_sub)

    return slide


def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets,
                         header_size=20, bullet_size=16):
    """Slide avec deux colonnes, taille d'en-tête configurable."""
    num = next_slide_num()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_red_bar(slide, top=0, height=Inches(1.2))
    add_bottom_bar(slide)
    add_page_number(slide, num)
    add_logo(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Colonne gauche
    txL = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
    tf_l = txL.text_frame
    p_l = tf_l.paragraphs[0]
    p_l.text = left_title
    p_l.font.size = Pt(header_size)
    p_l.font.bold = True
    p_l.font.color.rgb = RED

    txLC = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.6))
    tf_lc = txLC.text_frame
    tf_lc.word_wrap = True
    for i, b in enumerate(left_bullets):
        _add_styled_bullet_col(tf_lc, b, i == 0, bullet_size)

    # Colonne droite
    txR = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5))
    tf_r = txR.text_frame
    p_r = tf_r.paragraphs[0]
    p_r.text = right_title
    p_r.font.size = Pt(header_size)
    p_r.font.bold = True
    p_r.font.color.rgb = RED

    txRC = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.6))
    tf_rc = txRC.text_frame
    tf_rc.word_wrap = True
    for i, b in enumerate(right_bullets):
        _add_styled_bullet_col(tf_rc, b, i == 0, bullet_size)

    return slide


def _add_zone_box(slide, left, top, width, height, label, color):
    """Zone rectangulaire avec label en haut."""
    zone = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    zone.fill.solid()
    zone.fill.fore_color.rgb = RGBColor(0xF9, 0xF9, 0xF9)
    zone.line.color.rgb = color
    zone.line.width = Pt(2)

    txZ = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.3))
    tf_z = txZ.text_frame
    p_z = tf_z.paragraphs[0]
    p_z.text = label
    p_z.font.size = Pt(12)
    p_z.font.bold = True
    p_z.font.color.rgb = color
    p_z.alignment = PP_ALIGN.LEFT


def _add_component_box(slide, left, top, width, height, label, sublabel, color):
    """Boîte de composant avec label et sous-label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if sublabel:
        p2 = tf.add_paragraph()
        p2.text = sublabel
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xFF, 0xDD, 0xDD)
        p2.alignment = PP_ALIGN.CENTER


def _add_logo_on_box(slide, logo_name, box_left, box_top, box_width, logo_h=Inches(0.28)):
    """Place un logo dans le coin haut droit d'une pastille."""
    left = box_left + box_width - logo_h - Inches(0.04)
    top = box_top + Inches(0.04)
    safe_add_picture(slide, logo_name, left, top, logo_h)


def _add_sub_group_box(slide, left, top, width, height, label, border_color,
                       bg_color=None, font_size=10):
    """Sous-groupe avec bordure et label en haut à gauche."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or RGBColor(0xF5, 0xF0, 0xF0)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

    tx = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(0.03),
        width - Inches(0.16), Inches(0.25)
    )
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = border_color
    p.alignment = PP_ALIGN.LEFT


def _add_arrow(slide, left, top, width, color):
    """Flèche horizontale."""
    arr = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.25)
    )
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def add_architecture_slide():
    """Slide architecture data avec zones, sous-groupes dbt, logos."""
    num = next_slide_num()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_red_bar(slide, top=0, height=Inches(1.2))
    add_bottom_bar(slide)
    add_page_number(slide, num)
    add_logo(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Architecture data"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ── ZONE 1 : SOURCES ──
    _add_zone_box(slide, Inches(0.3), Inches(1.5), Inches(1.7), Inches(5.2), "SOURCES", BLUE)

    _add_component_box(slide, Inches(0.45), Inches(2.1), Inches(1.4), Inches(0.9),
                       "MySQL RDS", "winstat", BLUE)
    _add_logo_on_box(slide, "mysql", Inches(0.45), Inches(2.1), Inches(1.4))

    _add_component_box(slide, Inches(0.45), Inches(3.2), Inches(1.4), Inches(0.65),
                       "4 tables CDC", "temps réel", BLUE)

    _add_component_box(slide, Inches(0.45), Inches(4.05), Inches(1.4), Inches(0.65),
                       "14 tables REF", "bulk quotidien", BLUE)

    safe_add_picture(slide, "aws", Inches(0.65), Inches(5.0), Inches(0.5))

    # Flèche Sources → Ingestion
    _add_arrow(slide, Inches(2.0), Inches(3.2), Inches(0.4), BLUE)

    # ── ZONE 2 : INGESTION ──
    _add_zone_box(slide, Inches(2.4), Inches(1.5), Inches(2.1), Inches(5.2), "INGESTION", DARK_RED)

    _add_component_box(slide, Inches(2.6), Inches(2.1), Inches(1.7), Inches(0.75),
                       "Debezium 2.7", "CDC binlog", DARK_RED)
    _add_logo_on_box(slide, "debezium", Inches(2.6), Inches(2.1), Inches(1.7))

    _add_component_box(slide, Inches(2.6), Inches(3.0), Inches(1.7), Inches(0.75),
                       "Kafka", "4 topics streaming", DARK_RED)
    _add_logo_on_box(slide, "kafka", Inches(2.6), Inches(3.0), Inches(1.7))

    _add_component_box(slide, Inches(2.6), Inches(3.9), Inches(1.7), Inches(0.75),
                       "Python CDC", "micro-batch 500", DARK_RED)
    _add_logo_on_box(slide, "python", Inches(2.6), Inches(3.9), Inches(1.7))

    _add_component_box(slide, Inches(2.6), Inches(4.8), Inches(1.7), Inches(0.75),
                       "Python Bulk", "Parquet + COPY INTO", DARK_RED)
    _add_logo_on_box(slide, "python", Inches(2.6), Inches(4.8), Inches(1.7))

    # Flèche Ingestion → Snowflake
    _add_arrow(slide, Inches(4.5), Inches(3.2), Inches(0.3), DARK_RED)

    # ── ZONE 3 : SNOWFLAKE DWH ──
    _add_zone_box(slide, Inches(4.8), Inches(1.5), Inches(5.4), Inches(5.2), "     SNOWFLAKE DWH", RED)
    # Logo Snowflake au niveau de la zone (pas au niveau RAW)
    safe_add_picture(slide, "snowflake", Inches(8.2), Inches(1.55), Inches(0.45))

    # RAW (directement dans Snowflake DWH, alimenté par Python — pas par dbt)
    _add_component_box(slide, Inches(5.0), Inches(2.7), Inches(1.1), Inches(1.5),
                       "RAW", "18 tables\nbrutes", RGBColor(0xE0, 0x40, 0x40))

    # Flèche RAW → Transformation dbt
    _add_arrow(slide, Inches(6.1), Inches(3.35), Inches(0.25), RED)

    # ── Sous-groupe Transformation dbt (STAGING + MARTS) ──
    _add_sub_group_box(slide, Inches(6.55), Inches(2.35), Inches(3.35), Inches(2.35),
                       "Transformation dbt", RGBColor(0xC0, 0x60, 0x60),
                       RGBColor(0xFD, 0xF8, 0xF8))
    safe_add_picture(slide, "dbt", Inches(8.15), Inches(2.37), Inches(0.22))

    # STAGING
    _add_component_box(slide, Inches(6.75), Inches(2.9), Inches(1.2), Inches(1.5),
                       "STAGING", "18 modèles\ndédup + PII", RED)

    # Flèche STAGING → MARTS
    _add_arrow(slide, Inches(7.95), Inches(3.5), Inches(0.2), RED)

    # ── Sous-groupe MARTS ──
    _add_sub_group_box(slide, Inches(8.2), Inches(2.65), Inches(1.55), Inches(2.0),
                       "MARTS", RED, RGBColor(0xFF, 0xF0, 0xF0), font_size=9)
    _add_component_box(slide, Inches(8.35), Inches(3.05), Inches(1.25), Inches(0.42),
                       "3 DIM", "", RED)
    _add_component_box(slide, Inches(8.35), Inches(3.52), Inches(1.25), Inches(0.42),
                       "8 FAITS", "", RED)
    _add_component_box(slide, Inches(8.35), Inches(3.99), Inches(1.25), Inches(0.42),
                       "15 KPIs", "", RED)

    # ── Sous-groupe dbt en bas (AUDIT + SNAPSHOTS) ──
    # Couvre la largeur de RAW à Transformation dbt sans les contenir
    _add_sub_group_box(slide, Inches(5.0), Inches(5.0), Inches(4.9), Inches(1.45),
                       "dbt", RGBColor(0xA0, 0x30, 0x30), RGBColor(0xFA, 0xF5, 0xF5))
    safe_add_picture(slide, "dbt", Inches(5.4), Inches(5.02), Inches(0.22))

    # AUDIT
    _add_component_box(slide, Inches(5.15), Inches(5.4), Inches(2.2), Inches(0.65),
                       "AUDIT", "lineage + runs", RGBColor(0xA0, 0x30, 0x30))

    # SNAPSHOTS
    _add_component_box(slide, Inches(7.55), Inches(5.4), Inches(2.2), Inches(0.65),
                       "SNAPSHOTS", "SCD2 dimensions", RGBColor(0xA0, 0x30, 0x30))

    # Flèche Snowflake → Exposition
    _add_arrow(slide, Inches(10.2), Inches(3.2), Inches(0.35), GREEN)

    # ── ZONE 4 : EXPOSITION ──
    _add_zone_box(slide, Inches(10.55), Inches(1.5), Inches(2.5), Inches(5.2), "EXPOSITION", GREEN)

    _add_component_box(slide, Inches(10.75), Inches(2.1), Inches(2.1), Inches(0.65),
                       "Power BI", "dashboards", GREEN)
    _add_logo_on_box(slide, "powerbi", Inches(10.75), Inches(2.1), Inches(2.1))

    _add_component_box(slide, Inches(10.75), Inches(2.95), Inches(2.1), Inches(0.65),
                       "Tableau", "visualisation", GREEN)
    _add_logo_on_box(slide, "tableau", Inches(10.75), Inches(2.95), Inches(2.1), Inches(0.22))

    _add_component_box(slide, Inches(10.75), Inches(3.8), Inches(2.1), Inches(0.65),
                       "Metabase", "self-service BI", GREEN)
    _add_logo_on_box(slide, "metabase", Inches(10.75), Inches(3.8), Inches(2.1))

    _add_component_box(slide, Inches(10.75), Inches(4.65), Inches(2.1), Inches(0.65),
                       "API / Exports", "intégrations", GREEN)

    # ── ZONE TRANSVERSE : ORCHESTRATION ──
    orch = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.45)
    )
    orch.fill.solid()
    orch.fill.fore_color.rgb = ORANGE
    orch.line.fill.background()
    tf_o = orch.text_frame
    tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_o = tf_o.paragraphs[0]
    p_o.text = ("ORCHESTRATION & MONITORING    |    batch_loop.sh (5-30 min)    "
                "|    Docker Compose (6 services)    |    CI/CD GitHub Actions → GHCR    "
                "|    Alertes Teams    |    Source freshness")
    p_o.font.size = Pt(11)
    p_o.font.bold = True
    p_o.font.color.rgb = WHITE
    p_o.alignment = PP_ALIGN.CENTER

    return slide


def add_notes(slide, text):
    """Ajoute des notes du présentateur au slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# =============================================================================
# GÉNÉRATION DES SLIDES
# =============================================================================

# Slide 1 : Couverture
s = add_title_slide(
    "MediCore",
    "Pipeline ELT industrialisé pour données pharmacie"
)
add_notes(s, "Bonjour à tous. Je vais vous présenter le projet MediCore, un pipeline de données "
"développé pour centraliser et exploiter les données pharmacie. "
"L'objectif : transformer des données brutes dispersées dans MySQL en indicateurs "
"métier exploitables dans Snowflake, le tout de manière automatisée et fiable.")

# Slide 2 : Sommaire
s = add_sommaire_slide([
    "1. Contexte et problématique",
    "2. Architecture du pipeline",
    "3. Ingestion des données",
    "4. Transformations dbt",
    "5. Opérations et monitoring",
    "6. CI/CD",
    "7. Bilan et feuille de route",
])
add_notes(s, "On va parcourir ensemble 7 sections.\n"
"On commence par le contexte métier — pourquoi ce pipeline existe — puis on plonge "
"dans l'architecture technique, l'ingestion des données, les transformations dbt, "
"les opérations et le monitoring, le CI/CD, et enfin le bilan avec la feuille de "
"route.")

# --- SECTION 1 ---
s = add_section_slide(1, "Contexte et problématique")
add_notes(s, "Commençons par comprendre le besoin métier qui a motivé la création de MediCore.")

s = add_content_slide("Le besoin", [
    "Centraliser : données de pharmacies dispersées dans MySQL (winstat)",
    "  - 18 tables sources : ventes, commandes, stock, produits, fournisseurs…",
    "  - 920M+ lignes à charger et maintenir à jour",
    "KPIs métier : marge, écoulement, ruptures, ABC Pareto, trésorerie…",
    "Fraîcheur : données transactionnelles mises à jour en continu (5-30 min)",
    "Sécurité PII : masquage obligatoire des noms, adresses, téléphones",
])
add_notes(s, "Le point de départ, c'est une base de données MySQL appelée winstat qui contient "
"toutes les données opérationnelles des pharmacies : ventes, commandes, stock, "
"produits, fournisseurs. 18 tables au total avec plus de 920 millions de lignes.\n\n"
"Le problème : ces données sont enfermées dans un système transactionnel, pas "
"exploitables pour des analyses. Les utilisateurs ont besoin de KPIs métier — marge, "
"écoulement, ruptures, classification ABC — et ces calculs nécessitent des "
"croisements entre tables qui seraient trop coûteux sur la base opérationnelle. "
"D'où le besoin d'un entrepôt dédié.\n\n"
"Pourquoi Snowflake ? Parce qu'il sépare le stockage du calcul, ce qui permet de "
"scaler indépendamment, et son architecture columnar est optimisée pour les "
"requêtes analytiques.\n\n"
"La fraîcheur est un enjeu majeur : les données transactionnelles doivent être "
"mises à jour toutes les 5 à 30 minutes. Ces valeurs sont par défaut pour les "
"environnement de DEV et la PROD. Un pharmacien qui consulte ses ventes du jour a "
"besoin de données quasi temps réel, pas de données de la veille. C'est ce qui "
"justifie le CDC temps réel (Change Data Capture).\n\n"
"Enfin, la sécurité PII (Données personnelles identifiables). Le RGPD nous impose "
"de masquer les données personnelles — noms de patients, adresses, téléphones. On "
"utilise un hash MD5 irréversible appliqué dès la couche staging, avant toute "
"exploitation analytique. Aucune donnée personnelle en clair ne passe dans les "
"marts ou les dashboards.")

s = add_content_slide("Les contraintes", [
    "Deux flux avec latences différentes :",
    "  - 4 tables transactionnelles CDC temps réel : COMMANDES, FACTURES, ORDERS, MODSTOCK",
    "  - 14 tables référentiel bulk quotidien : PRODUITS, PHARMACIE, FOURNISSEURS…",
    "Volumétrie : chunks de 500K lignes, conteneur limité à 7.5 GB RAM",
    "Fiabilité : aucune perte de données tolérée, dead-letter queue pour les erreurs",
    "Traçabilité : chaque donnée reliée au batch qui l'a produite (lineage opérationnel)",
])
add_notes(s, "La difficulté principale, c'est qu'on a deux types de données très différents. "
"D'un côté, on a 4 tables transactionnelles qui changent en permanence — COMMANDES, "
"FACTURES, ORDERS, MODSTOCK — et elles nécessitent un suivi quasi temps réel via "
"CDC. De l'autre, on a 14 tables de référence — PRODUITS, PHARMACIE, "
"FOURNISSEURS… — ces tables sont plus stables et sont rechargées quotidiennement.\n\n"
"Cette dualité impose deux stratégies d'ingestion, chacune avec ses propres défis.\n"
"À cela s'ajoutent des contraintes :\n"
"- mémoire, le conteneur Docker est limité à 7.5 GO de RAM pour traiter 920 "
"millions de lignes\n"
"- des exigences de traçabilité : on doit pouvoir relier chaque donnée au batch "
"qui l'a produite\n"
"- et de fiabilité optimale avec aucune perte de données tolérée via le mécanisme "
"de DLQ (dead-letter queue) pour les erreurs.")

# --- SECTION 2 ---
s = add_section_slide(2, "Architecture du pipeline")
add_notes(s, "Voyons comment ces contraintes se traduisent concrètement en choix d'architecture.")

s = add_architecture_slide()
add_notes(s, "Voici l'architecture complète de Medicore.\n"
"De gauche à droite :\n"
"Les sources MySQL sur AWS RDS. Puis l'ingestion comprenant : Debezium qui capture "
"les changements du binlog MySQL, les publie sur de topics Kafka, Kafka les consomme "
"et les met à disposition aux deux scripts Python qui les insèrent dans Snowflake.\n"
"Pourquoi Kafka plutôt qu'une connexion directe ? Il y a 3 raisons à cela : le "
"découplage :\n"
"- si Snowflake est indisponible, Kafka conserve les messages;\n"
"- la possibilité de rejouer un offset en cas d'erreur;\n"
"- et la scalabilité qui permet de mettre plusieurs consumers en parallèle.\n\n"
"Dans Snowflake, RAW stocke les données brutes sans transformation (couche BRONZE) "
"— c'est le principe ELT.\n"
"Puis dbt transforme les données : les modèles staging (ou SILVER) dédupliquent et "
"masquent les données PII\n"
"Les modèles marts (ou GOLD) construisent le star schema pour donner en sortie "
"3 DIM, 8 FAITS et 15 KPIs.\n\n"
"L'AUDIT et les SNAPSHOTS sont deux composants essentiels des modèles dbt. L'AUDIT "
"comprenant les tables PIPELINE_RUNS, STEP_RUNS et DBT_MODEL_RUNS — permet la "
"traçabilité à chaque exécution du pipeline. Quand un KPI semble faux, on remonte "
"au RUN_ID, on vérifie chaque phase, et on identifie si c'est un problème "
"d'ingestion ou de transformation. Sans audit, le diagnostic prend des heures au "
"lieu de minutes.\n\n"
"Les SNAPSHOTS implémentent le SCD2 — Slowly Changing Dimension Type 2. Quand une "
"dimension change — exemple : une pharmacie qui change de nom, un produit qui change "
"de fournisseur — au lieu d'écraser l'ancienne valeur, on la conserve avec une date "
"de fin et on crée une nouvelle ligne avec la nouvelle valeur et une date de début. "
"Ça permet des analyses historiques cohérentes : les ventes de février s'affichent "
"avec l'ancien nom, celles de mars avec le nouveau.\n\n"
"À droite, la couche exposition. Trois outils sont envisagés :\n"
"- Power BI : intégration Microsoft native, avec licence souvent incluse dans "
"Microsoft 365, DAX, un langage de manipulation de données puissant.\n"
"Inconvénient : verrouillé dans l'écosystème Microsoft.\n"
"- Tableau : référence en data visualisation, excellente communauté, connexion "
"native Snowflake.\n"
"Inconvénient : son coût élevé (~70$/utilisateur/mois).\n"
"- Metabase : open-source, self-hosted, SQL natif, gratuit.\n"
"Inconvénient : moins de fonctionnalités avancées. Idéal pour démarrer ou du "
"self-service BI.\n"
"Le choix dépendra du budget, des besoins utilisateurs et de l'écosystème "
"existant.\n\n"
"L'API/Exports permet d'intégrer les KPIs dans d'autres systèmes — ERP, CRM, "
"applications métier. Les exports CSV/Excel permettent aussi de partager des données "
"avec des partenaires comme les laboratoires ou les groupements de pharmacies.\n\n\n"
"En bas, l'orchestration et le monitoring. Ils comprennent le script batch_loop.sh "
"qui orchestre les 7 phases séquentiellement toutes les 5 à 30 minutes. Docker "
"Compose gère les 6 services infrastructure. GitHub Actions automatise le CI/CD à "
"chaque push vers la branche main. Les alertes Teams notifient l'équipe en temps "
"réel.\n\n"
"GHCR — GitHub Container Registry — est le registre Docker intégré à GitHub. Après "
"chaque merge sur main et passage des 4 jobs CI, l'image Docker est poussée sur "
"ghcr.io/auganmadet/medicore avec un tag SHA pour la traçabilité exacte et un tag "
"« latest » pour la dernière version stable.\n\n"
"Les mécanismes Alertes Teams et Sources Freshness seront détaillés un peu plus "
"loin.\n"
"Mais retenez que Alertes Teams est mis en place via un webhook Microsoft (Slide 11) "
"et que Freshness (slide 16) est un mécanisme dbt pour contrôler la fraicheur des "
"données.")

s = add_content_slide("Infrastructure Docker", [
    "6 services Docker Compose :",
    "  - medicore_elt_batch : conteneur principal (CDC + dbt + orchestration)",
    "  - mysql_cdc : base de données source (MySQL 8.0)",
    "  - kafka + zookeeper : streaming des événements CDC",
    "  - connect : Debezium 2.7.3 (capture des changements MySQL)",
    "  - kafdrop : interface web de monitoring Kafka",
    "Sécurité : ports bindés sur 127.0.0.1, conteneur non-root (appuser)",
    "Ressources : 8 GB RAM, 2 CPUs pour le conteneur principal",
])
add_notes(s, "Pourquoi Docker Compose et pas Kubernetes ? Parce que notre cas d'usage est "
"simple : un pipeline batch mono-instance. Docker Compose nous donne un environnement "
"reproductible avec 6 services sur un réseau interne, sans la complexité "
"opérationnelle de Kubernetes.\n\n"
"Le rôle de chaque service :\n"
"- medicore_elt_batch : le cerveau du pipeline. Il exécute le batch_loop.sh, les "
"scripts Python d'ingestion CDC et bulk load, et toutes les commandes dbt.\n"
"- mysql_cdc : l'instance MySQL source contenant les données pharmacie winstat. "
"Le binlog est activé pour permettre le CDC.\n"
"- kafka : le broker de messages qui reçoit les événements CDC de Debezium et les "
"met à disposition des consommateurs. Garantit la durabilité et le replay des "
"messages.\n"
"- zookeeper : le gestionnaire de cluster Kafka. Il maintient la configuration, "
"les leaders de partitions et les métadonnées. Pourquoi l'associer à Kafka ? Parce "
"que Kafka en version 7.5.0 dépend encore de Zookeeper pour la coordination "
"distribuée. Les versions récentes avec KRaft s'en affranchissent, mais on utilise "
"la version stable.\n"
"- connect : Kafka Connect avec le connecteur Debezium 2.7.3. C'est le pont entre "
"MySQL et Kafka : il lit le binlog MySQL et publie les changements sur les topics.\n"
"- kafdrop : interface web pour visualiser les topics, les messages, les offsets et "
"le lag des consommateurs. Indispensable pour le diagnostic.\n\n"
"Les ports sont bindés sur 127.0.0.1 — pas d'exposition externe — et le conteneur "
"principal tourne avec un utilisateur non-root (appuser) pour limiter la surface "
"d'attaque.\n\n"
"Pourquoi 8 GB de RAM ? Le bulk load de 920M lignes par chunks de 500K nécessite "
"de la mémoire pour les DataFrames Parquet, plus la JVM Kafka/Connect (~2 GB), plus "
"dbt. 8 GB est le minimum stable observé. Pourquoi 2 CPUs ? Un pour "
"le pipeline Python/dbt, un pour Kafka/Connect. C'est suffisant pour un traitement "
"séquentiel mono-instance.")

# --- SECTION 3 ---
s = add_section_slide(3, "Ingestion des données")
add_notes(s, "Passons au cœur du système : comment les données transitent de MySQL à Snowflake.")

s = add_two_column_slide(
    "Deux flux d'ingestion",
    "CDC temps réel (4 tables)",
    [
        "Debezium : lecture du binlog MySQL",
        "Kafka : publication sur 4 topics",
        "Consumer Python : micro-batch 500 events",
        "INSERT : écriture dans Snowflake RAW",
        "Commit Kafka : après flush réussi uniquement",
        "DLQ : events malformés conservés",
        "Cycle : toutes les 5-30 min",
    ],
    "Bulk load quotidien (14 tables)",
    [
        "MySQL SELECT : cursor streaming server-side",
        "Parquet : écriture par chunks 500K rows",
        "Snowflake : PUT vers stage interne",
        "COPY INTO : avec FORCE=TRUE",
        "Rechargement : TRUNCATE + reload à 03h00",
        "Mémoire : gc.collect() après chaque chunk",
        "Volume : 920M+ lignes chargées au total",
    ],
)
add_notes(s, "Les deux colonnes reflètent les deux stratégies d'ingestion.\n\n"
"Le workflow CDC temps réel, étape par étape :\n"
"1. Debezium lit le binlog MySQL en continu et détecte chaque INSERT, UPDATE et "
"DELETE sur les 4 tables transactionnelles: COMMANDES, FACTURES, ORDERS, MODSTOCK.\n"
"2. Il publie un événement JSON sur un topic Kafka dédié — un topic par table.\n"
"3. Notre consumer Python s'abonne aux 4 topics et accumule les messages.\n"
"4. Tous les 500 messages ou 30 secondes (délai timeout), il fait un flush : INSERT "
"batch dans Snowflake RAW.\n"
"5. Après confirmation de l'INSERT réussi, il commit l'offset Kafka manuellement.\n"
"6. Les messages malformés partent en DLQ.\n"
"Ce worflow s'inscrit dans une boucle tous les 5 min en DEV et 30 min en PROD\n\n"
"Pourquoi le commit Kafka manuel et pas l'auto-commit ? Parce que l'auto-commit "
"avance l'offset dès la lecture du message. Si l'INSERT échoue, le message est "
"perdu définitivement. Le commit manuel garantit le « at-least-once delivery » "
"cad la livraison du message.\n\n"
"Le workflow bulk load quotidien se déroule en plusieurs étapes :\n"
"1. Connexion MySQL avec cursor streaming server-side — pas pd.read_sql() qui "
"bufférise tout en mémoire, c'est un piège classique.\n"
"2. Lecture par chunks de 500K lignes depuis chaque table.\n"
"3. Conversion en format Parquet qui est un format columnar compressé, optimal "
"pour Snowflake.\n"
"4. PUT vers un stage interne Snowflake ce qui correspond à un upload du fichier "
"Parquet.\n"
"5. COPY INTO pour charger le stage dans la table RAW avec FORCE=TRUE.\n"
"6. gc.collect() + del df pour libérer la mémoire après chaque chunk.\n"
"Ce worflow permet de traiter 920M de lignes\n\n"
"Le FORCE=TRUE est critique : je me suis aperçu qu'après un TRUNCATE, Snowflake "
"garde en cache les métadonnées des fichiers pendant 64 jours et les ignore "
"silencieusement. Sans FORCE=TRUE, les données semblent chargées mais ne le sont "
"pas en vérité. C'est un piège qui m'a coûté du temps à diagnostiquer.")

s = add_content_slide("Fiabilité de l'ingestion", [
    "Dead-Letter Queue (DLQ) : messages CDC invalides conservés, jamais perdus",
    "Commit Kafka manuel : l'offset n'avance qu'après insertion Snowflake réussie",
    "Fallback row-by-row : si un batch INSERT échoue, insertion ligne par ligne",
    "Monitoring volume : alerte après N batches consécutifs à 0 events",
    "FORCE=TRUE : contourne le cache metadata Snowflake (persistant 64 jours)",
    "Alertes Teams : notification après 3 échecs consécutifs + recovery automatique",
])
add_notes(s, "La fiabilité, c'est ce qui différencie un pipeline de production d'un script "
"jetable. Je vais détailler chaque mécanisme.\n\n"
"Dead-Letter Queue (DLQ) : quand un message CDC est malformé (ex:/ un JSON invalide, "
"des champs manquants, un type incompatible), plutôt que de crasher ou d'ignorer, on "
"l'insère dans une table _DLQ avec le message d'erreur, le timestamp et le contenu "
"brut. On peut ensuite analyser, corriger la cause et éventuellement rejouer.\n"
"Les messages CDC invalides sont donc conservés et jamais perdus\n\n"
"Commit Kafka manuel : par défaut, Kafka auto-commit l'offset toutes les 5 secondes. "
"Le consumer lit un message et l'offset avance. Si l'INSERT Snowflake échoue ensuite, "
"le message est perdu — l'offset a déjà avancé. Le commit manuel inverse la logique : "
"pas d'INSERT réussi, pas d'avance d'offset. C'est le « at-least-once delivery ».\n\n"
"Fallback row-by-row : un batch INSERT de 500 lignes échoue si UNE seule ligne a un "
"problème (par ex:/ un type incompatible, une valeur trop longue). Le fallback "
"réessaie ligne par ligne : les 499 valides passent, la ligne problématique va en "
"DLQ avec le détail de l'erreur.\n\n"
"Monitoring volume : on compte les événements CDC par batch. Si N batches consécutifs "
"retournent 0 événements, on alerte via Teams. — Un topic Kafka vide n'est pas "
"toujours normal : Debezium peut avoir décroché, le binlog peut être purgé, ou un "
"problème réseau.\n"
"Le nombre de batches consécutifs est configurable via une variable d'environnement.\n\n"
"FORCE=TRUE : Snowflake maintient un registre des fichiers déjà chargés par COPY "
"INTO — ce cache persiste 64 jours même après un TRUNCATE. Sans FORCE=TRUE, les "
"mêmes noms de fichiers sont silencieusement ignorés lors du rechargement.\n\n"
"Alertes Teams : un webhook Microsoft Teams envoie des notifications en temps réel "
"à l'équipe. Le mécanisme fonctionne avec un compteur d'échecs par phase. Après 3 "
"échecs consécutifs sur une même phase (CDC, staging, marts...), une alerte critique "
"est envoyée en format Adaptive Card JSON. Quand le compteur repasse à 0 cad la "
"phase réussit à nouveau, une notification de recovery est envoyée automatiquement.\n"
"Un mécanisme de retry avec backoff exponentiel évite de perdre des alertes si Teams "
"est temporairement indisponible. L'URL du webhook est optionnelle : si elle n'est "
"pas configurée, le pipeline continue normalement en loggant sans alerter.")

# --- SECTION 4 ---
s = add_section_slide(4, "Transformations dbt")
add_notes(s, "Les données sont dans Snowflake RAW. Voyons comment dbt les transforme en "
"informations exploitables.")

s = add_content_slide("Staging : nettoyage et sécurité", [
    "18 modèles staging : un stg_*.sql par table source",
    "Déduplication CDC : ROW_NUMBER() OVER (PARTITION BY PK ORDER BY cdc_timestamp DESC)",
    "Filtre des deletes : WHERE cdc_operation != 'D'",
    "Masquage PII : hash MD5 sur noms, adresses, téléphones, emails",
    "  - Macro centralisée {{ mask_pii('colonne') }}",
    "  - Jamais de données personnelles en clair dans STAGING ou MARTS",
    "Matérialisation incrémentale : merge sur clé primaire composite",
])
add_notes(s, "Pourquoi dbt plutôt que des procedures stockées Snowflake pour la transformation "
"des données ?\n"
"Trois avantages majeurs à choisir dbt :\n"
"- le versioning Git de chaque transformation SQL,\n"
"- les tests automatiques intégrés,\n"
"- et le lineage qui permet de savoir exactement d'où provient chaque colonne.\n\n"
"Les modèles dbt Staging s'appuyent sur les données sources de la couche RAW.\n"
"Le staging fait deux choses critiques. D'abord la déduplication CDC : avec "
"ROW_NUMBER partitionné par clé primaire et trié par timestamp décroissant, on ne "
"garde que la version la plus récente de chaque ligne. Ensuite le masquage PII "
"avec une macro MD5 centralisée.\n\n"
"Le CDC capture toutes les opérations MySQL : les creates (c), updates (u), "
"deletes (d) et snapshots (r). Elles sont toutes stockées dans RAW pour traçabilité "
"complète.\n\n"
"Mais dans staging, on filtre les deletes parce que dans le contexte pharmacie, une "
"suppression dans MySQL ne signifie pas qu'on veut perdre la donnée analytiquement. "
"Exemples concrets :\n"
"- Une commande archivée et supprimée de MySQL reste nécessaire pour calculer le "
"CA annuel\n"
"- Une facture purgée de la base opérationnelle doit rester dans l'entrepôt pour "
"les analyses de marge historiques\n"
"- Un mouvement de stock supprimé après clôture doit rester pour les rapports "
"d'inventaire\n\n"
"Le staging ne garde que l'état courant \"vivant\" de chaque ligne — la dernière "
"version non supprimée. Si un analyste a besoin de savoir qu'une ligne a été "
"supprimée (et quand), il consulte RAW où le cdc_operation = 'D' avec son "
"cdc_timestamp est conservé.\n\n"
"Qu'est-ce que la matérialisation incrémentale ? Par défaut, dbt recrée entièrement "
"chaque table à chaque exécution — c'est le mode « table ». Ça fonctionne pour les "
"petits volumes, mais avec 920 millions de lignes, recalculer tout à chaque cycle "
"de 5 minutes serait intenable. La matérialisation incrémentale ne traite que les "
"nouvelles données : elle filtre avec cdc_timestamp >= max(loaded_at), donc seules "
"les lignes arrivées depuis le dernier run sont transformées. La stratégie « merge » "
"fusionne ces lignes avec l'existant via la clé primaire composite — si la ligne "
"existe déjà, elle est mise à jour ; sinon, elle est insérée. Résultat : un run "
"staging prend quelques secondes au lieu de dizaines de minutes.\n\n"
"Résultat : Les modèles Staging crééent les tables de la couche SILVER.")

s = add_content_slide("Marts : star schema et KPIs", [
    "3 dimensions : dim_pharmacie, dim_produit, dim_fournisseur (surrogate keys)",
    "8 tables de faits :",
    "  - fact_ventes, fact_commandes, fact_prix_journalier, fact_stock_mouvement",
    "  - fact_stock_valorisation, fact_ruptures, fact_tresorerie, fact_operateur",
    "15 KPIs métier :",
    "  - Marge, écoulement, ABC (Pareto), stock, stock valorisation",
    "  - Ruptures, trésorerie, opérateur, qualité des données",
    "Snapshots SCD2 : historisation pharmacie, produit, fournisseur",
])
add_notes(s, "Les modèles Marts s'appuyent sur les données de la couche SILVER\n"
"Pourquoi choisir un modèle star schema et pas un data vault ? Parce que notre "
"besoin est la performance analytique, pas la flexibilité d'intégration "
"multi-sources. Le star schema est compris par tous les outils BI, les jointures "
"sont simples et performantes.\n\n"
"Les modèles Marts crééent les tables de Dimension, les tables de Faits et des "
"tables Kpis.\n\n"
"3 dimensions créées :\n"
"- dim_pharmacie : identifie chaque pharmacie (nom, adresse, ville, département). "
"Clé métier : PHA_ID.\n"
"- dim_produit : représente chaque produit pharmaceutique enrichi avec les codes "
"EAN13 (code-barres) et LPPR (remboursement).\n"
"    Clé métier : PRD_ID.\n"
"- dim_fournisseur : représentent les répartiteurs et laboratoires qui fournissent "
"les produits. Clé métier : FOU_ID.\n\n"
"Chaque dimension utilise une clé de substitution qu'on appelle Surrogate keys "
"générée par ROW_NUMBER() — un entier séquentiel (1, 2, 3...) au lieu de la clé "
"métier source.\n"
"Pourquoi ne pas garder directement PHA_ID ou PRD_ID ? Trois raisons :\n"
"1. Indépendance vis-à-vis du système source — si MySQL change son schéma de clés, "
"les marts ne sont pas impactés\n"
"2. Performance — les jointures sur des entiers sont plus rapides que sur des "
"chaînes\n"
"3. Gestion des orphelins — quand une fact référence un produit ou une pharmacie "
"qui n'existe pas dans la dimension, on attribue la surrogate key -1 via "
"COALESCE(dim.sk, -1) au lieu de perdre la ligne\n\n"
"8 tables de faits créées:\n"
"1. fact_ventes : ventes journalières par pharmacie/produit (quantité, CA, marge)\n"
"2. fact_commandes : commandes passées aux fournisseurs (quantité, montant)\n"
"3. fact_prix_journalier : évolution quotidienne des prix d'achat et de vente\n"
"4. fact_stock_mouvement : mouvements de stock (entrées, sorties, ajustements)\n"
"5. fact_stock_valorisation : valorisation du stock à un instant T (quantité × "
"prix)\n"
"6. fact_ruptures : ruptures de stock avec estimation du CA perdu\n"
"7. fact_tresorerie : flux de trésorerie (encaissements, décaissements)\n"
"8. fact_operateur : performance des vendeurs (nombre de ventes, panier moyen)\n\n"
"Chaque fait référence les 3 dimensions via LEFT JOIN + COALESCE(sk, -1) pour ne "
"jamais perdre de lignes même si la dimension est incomplète.\n\n"
"Les 15 KPIs couvrent les principaux axes d'analyse. Notamment le KPI opérateur "
"(fact_operateur) qui mesure la performance de chaque vendeur en pharmacie via le "
"nombre de ventes, panier moyen, CA généré. C'est un levier managérial concret : "
"le titulaire peut identifier ses meilleurs vendeurs et adapter ses formations.\n\n"
"La classification ABC applique la loi de Pareto : elle identifie les 20%% de "
"produits générant 80%% du CA (classe A), les 30%% suivants pour 15%% du CA "
"(classe B), et les 50%% restants pour 5%% du CA (classe C). Ça guide les décisions "
"de réapprovisionnement et de négociation fournisseur.\n\n"
"Le KPI qualité des données surveille la fraîcheur et les anomalies — c'est le "
"monitoring du pipeline lui-même, il est vu depuis les marts.\n\n"
"Pourquoi historiser uniquement pharmacie, produit et fournisseur en SCD2, et pas "
"les 14 autres tables ? Parce que ce sont les seules dimensions dont les changements "
"ont un impact analytique. Si une pharmacie change de nom — par exemple « Pharmacie "
"Dupont » devient « Pharmacie Santé Plus » après un rachat — les ventes de janvier "
"doivent rester associées à l'ancien nom, et celles de février au nouveau. Sans "
"SCD2, toutes les ventes apparaîtraient sous « Santé Plus », cela fausse les "
"analyses.\n\n"
"Autre exemple concret : un produit change de fournisseur. Le Doliprane 500mg était "
"fourni par le répartiteur Alliance Healthcare, mais passe chez OCP. Le SCD2 conserve "
"les deux lignes avec des dates de validité : dbt_valid_from et dbt_valid_to. Les "
"analyses de marge par fournisseur restent correctes sur chaque période.\n\n"
"Les autres tables — EAN13, LPPR, DAYBYDAY — sont des référentiels stables ou des "
"données factuelles qui ne changent pas rétroactivement. Historiser chaque table "
"augmenterait la complexité sans apport analytique.")

# --- SECTION 5 ---
s = add_section_slide(5, "Opérations et monitoring")
add_notes(s, "Un pipeline n'a de valeur que s'il est fiable en production. Voyons les "
"mécanismes de supervision.")

s = add_content_slide("Orchestration et audit", [
    "batch_loop.sh : boucle principale (5 min dev / 30 min prod)",
    "  - 7 phases : ref_reload, cdc_batch, dbt_staging, dbt_snapshot, dbt_marts, dbt_test, freshness",
    "  - Timeout configurable par phase (défaut 30 min)",
    "  - Arrêt graceful sur SIGTERM/SIGINT (compatible docker compose stop)",
    "RUN_ID unique : UUID généré par itération de batch",
    "  - Schéma AUDIT Snowflake : PIPELINE_RUNS, PIPELINE_STEP_RUNS, DBT_MODEL_RUNS",
    "  - Chaque donnée produite est reliée au batch qui l'a générée",
    "Source freshness : CDC 12h warn / 24h error, référence 36h warn / 48h error",
])
add_notes(s, "Le batch_loop.sh est volontairement simple — un script bash plutôt qu'Airflow "
"— parce que l'orchestration est linéaire : 7 phases séquentielles. Pas besoin d'un "
"DAG complexe pour l'instant.\n\n"
"Les 7 phases dans l'ordre :\n"
"1. ref_reload — rechargement quotidien des 14 tables de référence à 03h00. TRUNCATE "
"+ bulk load complet. Ne s'exécute qu'une fois par jour.\n"
"2. cdc_batch — consommation des événements Kafka sur les 4 tables CDC. Micro-batch "
"de 500 events ou timeout 30s. C'est la phase la plus fréquente.\n"
"3. dbt_staging — transformation des données brutes RAW en staging : déduplication "
"CDC, masquage PII, typage. Run incrémental sur les 18 modèles stg_*.\n"
"4. dbt_snapshot — mise à jour des snapshots SCD2 : pharmacie, produit, fournisseur. "
"Détecte les changements et historise.\n"
"5. dbt_marts — construction du star schema : dimensions, faits, KPIs. Les marts "
"sont la couche consommée par les outils BI.\n"
"6. dbt_test — exécution des tests dbt : not_null, unique, relationships, "
"expression_is_true sur les formules métier. Les échecs sont comptabilisés.\n"
"7. freshness — contrôle de fraîcheur dbt source freshness sur les 18 tables RAW.\n\n"
"Le timeout est configurable par phase — PHASE_TIMEOUT_SEC, défaut 30 minutes — "
"est un retour d'expérience critique. Sans timeout, un dbt run bloqué sur une "
"requête Snowflake pouvait bloquer tout le pipeline indéfiniment. En production, "
"une requête sur une table mal clusterisée a déjà pris 45 minutes. Avec le timeout, "
"après 30 minutes on coupe la phase, on log l'échec dans AUDIT, on alerte Teams, "
"et le pipeline continue avec la phase suivante.\n\n"
"L'arrêt graceful intercepte les signaux SIGTERM et SIGINT. Quand on fait docker "
"compose stop, Docker envoie SIGTERM au processus. Sans gestion de ce signal, le "
"pipeline serait tué en plein milieu d'un INSERT ou d'un dbt run, on risque d'avoir "
"des données corrompues. Avec le trap, le pipeline termine proprement la phase en "
"cours avant de s'arrêter. C'est compatible avec Docker.\n\n"
"Chaque itération génère un RUN_ID unique — un UUID — propagé à toutes les phases "
"et persisté dans le schéma AUDIT Snowflake. Concrètement, on peut répondre à la "
"question : « cette donnée dans les marts, elle vient de quel batch, à quelle heure, "
"et est-ce que toutes les phases ont réussi ? »\n\n"
"Le source freshness vérifie que les données sont à jour. Deux seuils différents "
"selon le type : les tables CDC doivent être fraîches à moins de 12h (warning) ou "
"24h (erreur), car elles sont alimentées toutes les 5-30 minutes. Les tables "
"référence ont des seuils plus larges — 36h warning, 48h erreur — car elles ne sont "
"rechargées qu'une fois par jour à 03h00. Si une table dépasse le seuil, l'équipe "
"est alertée via Teams.")

# --- SECTION 6 ---
s = add_section_slide(6, "CI/CD")
add_notes(s, "Un pipeline sans CI/CD, c'est un pipeline qui casse en silence. Voyons les "
"garde-fous automatisés.")

s = add_content_slide("Pipeline CI/CD GitHub Actions", [
    "4 jobs CI en parallèle à chaque push :",
    "  - Lint Python (flake8) : imports, syntaxe, formatage",
    "  - Validation dbt (dbt parse) : syntaxe SQL/Jinja2, ref() et source() valides",
    "  - Build Docker : vérification que l'image se construit",
    "  - ShellCheck : analyse statique des scripts bash",
    "1 job CD conditionnel :",
    "  - Push image Docker vers GitHub Container Registry (GHCR)",
    "  - Déclenchement : uniquement sur main, après passage des 4 jobs CI",
    "  - Tags : SHA du commit + latest",
    "  - Authentification : automatique via GITHUB_TOKEN",
])
add_notes(s, "4 jobs CI tournent en parallèle en 2 à 3 minutes:\n"
"- Le lint Python attrape les imports inutilisés et les problèmes de formatage.\n"
"- Le dbt parse valide la syntaxe SQL et Jinja2 sans connexion Snowflake — on "
"utilise des variables factices.\n"
"- Docker build vérifie l'image.\n"
"- ShellCheck analyse les scripts bash.\n\n"
"Si les 4 jobs passent et qu'on est sur main, le CD pousse l'image sur GHCR. "
"Pourquoi GHCR plutôt que Docker Hub ? Parce que c'est intégré nativement à GitHub, "
"l'authentification est automatique via GITHUB_TOKEN (zéro secret à configurer "
"donc) et c'est gratuit pour les repos publics.\n\n"
"Le CD conditionnel concerne les tags Docker qui constituent un point important. "
"Chaque image poussée reçoit deux tags :\n"
"- Le tag SHA du commit — par exemple « a1b2c3d » — c'est l'identifiant unique et "
"immuable. Il permet de savoir exactement quel code tourne en production. Si un bug "
"apparaît, on remonte au commit exact, on lit le diff, on identifie la cause. C'est "
"la traçabilité totale du code au conteneur.\n"
"- Le tag « latest » — c'est un alias mobile qui pointe toujours sur la dernière "
"image construite depuis main. Il permet de déployer simplement avec docker pull "
"sans connaître le SHA.\n\n"
"En production, on référence le SHA pour la reproductibilité ; en dev, on utilise "
"latest pour toujours avoir la dernière version.\n"
"Cette double stratégie combine le meilleur des deux mondes : traçabilité exacte "
"(SHA) et commodité opérationnelle (latest).")

# --- SECTION 7 ---
s = add_section_slide(7, "Bilan et feuille de route")
add_notes(s, "Pour conclure, faisons le bilan de ce qui est opérationnel et ce qui reste à faire.")

s = add_two_column_slide(
    "Bilan et feuille de route",
    "Implémenté et opérationnel",
    [
        "Pipeline ELT complet : 18 tables ingérées",
        "Star schema : 3 dim, 8 faits, 15 KPIs",
        "CDC temps réel : Debezium/Kafka (4 tables)",
        "Bulk quotidien : 14 tables référentiel",
        "Masquage PII : automatique via macro dbt",
        "Lineage : RUN_ID + schéma AUDIT Snowflake",
        "CI/CD : GitHub Actions → GHCR automatisé",
        "Monitoring : Teams + source freshness dbt",
    ],
    "Feuille de route (à implémenter)",
    [
        "Exposition BI : Power BI / Metabase sur les marts",
        "Monitoring Kafka : détection retard offset lag",
        "Tests isolés : dbt seeds/fixtures en CI",
        "Documentation : KPIs avec formules et exemples SQL",
        "Credentials : rotation automatique des 4 comptes",
        "Orchestration : Airflow ou Dagster (vs bash)",
    ],
    header_size=24,
    bullet_size=15,
)
add_notes(s, "À gauche, tout ce qui est implémenté et fonctionnel. Le pipeline complet "
"ingère les 18 tables, le star schema produit 3 dim, 8 faits et 15 KPIs, le CDC "
"tourne en temps réel, le batch Bulk load tourne quotidiennement, le masquage PII "
"est automatique, le lineage audit trace chaque exécution, et le CI/CD déploie "
"automatiquement.\n\n"
"À droite, ce sont les prochaines étapes. Détaillons chaque point :\n\n"
"Exposition BI : les marts sont prêts, il manque l'outil de visualisation. Power BI, "
"Tableau ou Metabase — le choix dépend du budget et de l'écosystème existant. C'est "
"la priorité numéro UNE car c'est ce qui rend les KPIs accessibles aux utilisateurs "
"(pharmaciens).\n\n"
"Monitoring Kafka : aujourd'hui on monitore le volume CDC — combien d'events par "
"batch — mais pas le « consumer lag », c'est-à-dire le retard entre le dernier "
"message publié par Debezium et le dernier message consommé par notre pipeline. Un "
"lag croissant signifie que le consumer ne suit pas le rythme de production. "
"L'idée est d'exposer cette métrique via Kafdrop pour déclencher des alertes avant "
"que le retard devienne critique.\n\n"
"Tests isolés : aujourd'hui, dbt test s'exécute uniquement dans le batch_loop — il "
"envoie des requêtes de vérification (not_null, unique, relationships) sur les "
"vraies données Snowflake. Ça fonctionne, mais on ne peut pas tester en CI car "
"GitHub Actions n'a pas de connexion Snowflake. L'objectif est de créer des modèles "
"« seeds » dbt cad des fichiers CSV avec quelques lignes de données fictives — que "
"dbt chargerait dans un schéma temporaire pour exécuter les tests à chaque pull "
"request, avant le merge, sans dépendre de l'environnement de production.\n\n"
"Documentation KPIs : chaque KPI doit être documenté avec sa formule exacte, un "
"exemple SQL concret et un cas d'usage métier. Par exemple : « Marge brute = "
"(prix_vente - prix_achat) / prix_vente × 100. Exemple : un Doliprane vendu 3.50€ "
"acheté 2.10€ = marge de 40%%. » Cette documentation est essentielle pour que les "
"utilisateurs comprennent et fassent confiance aux chiffres.\n\n"
"Rotation des credentials : les 4 comptes — Snowflake, MySQL, Kafka, Teams webhook — "
"ont des mots de passe statiques. L'objectif est d'automatiser leur rotation "
"périodique, idéalement via un gestionnaire de secrets comme AWS Secrets Manager ou "
"HashiCorp Vault. Aujourd'hui, la rotation peut être manuelle, ce qui représente un "
"risque de sécurité en cas de compromission.")

# Slide finale
num = next_slide_num()
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RED

pic = slide.shapes.add_picture(LOGO_PATH, Inches(0), Inches(0.8), height=Inches(1.2))
pic.left = int((prs.slide_width - pic.width) / 2)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Merci"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Questions ?"
p2.font.size = Pt(28)
p2.font.color.rgb = LIGHT_RED
p2.alignment = PP_ALIGN.CENTER

add_page_number(slide, num)
add_notes(slide, "J'en ai fini avec la présentation.\n"
"Merci pour votre attention. Je suis disponible pour répondre à vos "
"questions, que ce soit sur les choix d'architecture, les détails techniques, "
"ou les prochaines étapes du projet.")

# Sauvegarder
output_path = "docs/MediCore_Presentation.pptx"
prs.save(output_path)
print(f"Présentation générée : {output_path} ({slide_counter} slides)")
